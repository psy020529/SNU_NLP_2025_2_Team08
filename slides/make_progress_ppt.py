#!/usr/bin/env python3
"""
Generate a 5–7 minute progress presentation PowerPoint based on the current plan.
Requires: python-pptx
Output: slides/progress_presentation.pptx
"""
from __future__ import annotations

from pathlib import Path
from typing import List

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except Exception as e:
    raise SystemExit("Please install python-pptx: pip install python-pptx") from e


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: List[str],
    notes: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_presentation(output: Path) -> None:
    prs = Presentation()

    # Slide 1 — Title/Motivation/Problem
    add_title_slide(
        prs,
        title="Steering LLM Outputs and User Detectability",
        subtitle="Progress Brief (5–7 minutes)",
    )

    add_bullet_slide(
        prs,
        title="Motivation & Problem",
        bullets=[
            "LLMs influence information trust and decisions",
            "Outputs can be steered via prompts/fine-tuning (prior work)",
            "Q1: Can system prompts/LoRA shift responses?",
            "Q2: Can users/detectors recognize steering?",
            "Scope: Neutral/synthetic issues; no real targeting",
        ],
        notes=(
            "We measure and detect steering signals; we do not build persuasion systems.\n"
            "The study connects to instruction tuning, prompt design, and evaluation."
        ),
    )

    # Slide 2 — Datasets
    add_bullet_slide(
        prs,
        title="Datasets & Initial Analysis",
        bullets=[
            "Internal: prompts.csv (3–5 topics × 20 neutral prompts)",
            "Public (licensed samples): open_bbq, apa22, agora_speech, skotapa",
            "Quick EDA: length, language mix, duplicates, CSV hygiene",
        ],
        notes=(
            "Start with small, licensed samples for fast iteration and compliance."
        ),
    )

    # Slide 3 — Preprocessing
    add_bullet_slide(
        prs,
        title="Preprocessing & Splits",
        bullets=[
            "Normalization: unicode/whitespace, sentence splitting (if needed)",
            "Splits: baseline vs manipulated (ID-based)",
            "JSONL schema: {id, topic, frame, prompt, response, meta}",
        ],
        notes=(
            "Clean schemas enable comparable metrics and detector generalization."
        ),
    )

    # Slide 4 — Implementation
    add_bullet_slide(
        prs,
        title="Baselines & Implementation",
        bullets=[
            "Code skeleton ready: scripts/src/outputs/notebooks",
            "Steering: System prompts (NEUTRAL/PRO/CON) + LoRA SFT (100–300 pairs)",
            "Detection: Rule-based metrics + simple classifier (LogReg/MLP)",
        ],
        notes=(
            "Prompt vs LoRA lets us compare fast vs structural steering effects."
        ),
    )

    # Slide 5 — Metrics
    add_bullet_slide(
        prs,
        title="Metrics & Displays",
        bullets=[
            "FramingScore: (pro-lex − con-lex)/length",
            "BalanceScore: 1 − |pro−con|/(pro+con+ε)",
            "Diversity: distinct-n, TTR; Detector: F1/AUC",
        ],
        notes=(
            "We show a table for means by frame, a boxplot, and a ROC curve."
        ),
    )

    # Slide 6 — Plan
    add_bullet_slide(
        prs,
        title="Future Plan & Milestones",
        bullets=[
            "W1: Baseline inference (≥20 prompts) + first metrics",
            "W2: LoRA SFT (100–300 pairs) + detector v1",
            "W3: Ablations + final plots/report",
        ],
        notes=(
            "Targets: d≥0.5, F1≥0.75, AUC≥0.80, distinct-2 drop ≤10pp."
        ),
    )

    # Slide 7 — Risks
    add_bullet_slide(
        prs,
        title="Risks & Mitigations",
        bullets=[
            "Data/license/domain shift → use minimal samples + synthetic prompts",
            "Resource: small models, short gens, LoRA r/α tuning",
            "Metrics sensitivity/overfitting → multi-features, cross-topic test",
            "Ethics: neutral scope; no targeting; transparency",
        ],
        notes=(
            "We maintain safe scope and have fallback routes for progress."
        ),
    )

    # Slide 8 — Presentation Quality
    add_bullet_slide(
        prs,
        title="Presentation Quality",
        bullets=[
            "Logical flow: Intro→Data→Methods→Progress→Metrics→Plan→Risks",
            "Visual clarity: ≥14pt, ≤3 columns, labeled axes/legends",
            "Team split: Motivation/Data; Impl/Metrics; Plan/Risks",
            "Backup: folder tree, JSONL schema, TODO screenshot",
        ],
        notes=(
            "We timebox ~45–50s per slide and rotate speakers smoothly."
        ),
    )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


if __name__ == "__main__":
    build_presentation(Path(__file__).parent / "progress_presentation.pptx")
